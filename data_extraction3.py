import gradio as gr
from langchain_community.document_loaders import UnstructuredPDFLoader
import pytesseract
from docx import Document
from pptx import Presentation
import os

# Path to the Tesseract executable
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def extract_text_from_pdf(pdf_path):
    pdf_loader = UnstructuredPDFLoader(pdf_path)
    documents = pdf_loader.load()
    extracted_text = "\n".join([doc.page_content for doc in documents])
    return extracted_text

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    extracted_text = "\n".join([para.text for para in doc.paragraphs])
    return extracted_text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    extracted_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
    return extracted_text

def extract_text(file):
    if file.name.endswith('.pdf'):
        return extract_text_from_pdf(file.name)
    elif file.name.endswith('.docx'):
        return extract_text_from_docx(file.name)
    elif file.name.endswith('.pptx'):
        return extract_text_from_pptx(file.name)
    else:
        return "Unsupported file format. Please upload a PDF, DOCX, or PPTX file."

def save_text_to_file(text):
    # Save the extracted text to a .txt file
    output_file = "extracted_data.txt"
    with open(output_file, "w", encoding="utf-8") as text_file:
        text_file.write(text)
    return output_file  # Return the file path for download

def main():
    with gr.Blocks() as demo:
        with gr.Row():
            with gr.Column():
                file_input = gr.File(label="Upload PDF, DOCX, or PPTX File")
                submit_button = gr.Button("Extract Text")

            with gr.Column():
                output_text = gr.Textbox(label="Extracted Text", lines=10)
                download_button = gr.File(label="Download Extracted Text")

        extracted_text = gr.State()  # To hold the extracted text

        def handle_extract(file):
            text = extract_text(file)
            return text, text  # Return the extracted text for both output and state

        def handle_download(text):
            if text:
                return save_text_to_file(text)  # Save and return the file path
            return None  # Return None if no text is available for download

        submit_button.click(
            fn=handle_extract,
            inputs=file_input,
            outputs=[output_text, extracted_text]  # Return output to both the textbox and state
        )

        download_button.change(
            fn=handle_download,
            inputs=extracted_text,
            outputs=download_button
        )

    demo.launch()

if __name__ == "__main__":
    main()
